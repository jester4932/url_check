import re

from pptx import Presentation


class PPTXRead:
    def __init__(self, **kwargs):
        self.file = kwargs.get("file", "")
        self.regex = kwargs.get("regex", "")
        self.text_frame = ""
        self.url = []

    def main(self):
        self.plain_text()
        self.hyperlinks()
        return self.url

    def plain_text(self):
        prs = Presentation(self.file)
        text_content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    self.text_frame = shape.text_frame
                    for paragraph in self.text_frame.paragraphs:
                        text_content.append(paragraph.text)
                        url = re.findall(self.regex, paragraph.text)
                        if url:
                            self.url.append(paragraph.text)

    def hyperlinks(self):
        for paragraph in self.text_frame.paragraphs:
            for run in paragraph.runs:  # find embedded hyperlinks
                if run.hyperlink.address:
                    self.url.append(run.hyperlink.address)
